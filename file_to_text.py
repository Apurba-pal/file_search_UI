# pip install python-docx openpyxl python-pptx pandas PyMuPDF

import os
import pandas as pd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import fitz 


def extract_text_from_word(docx_path):
    print(f"\n[Extracting text from Word]: {docx_path}")
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            print(para.text)
    except Exception as e:
        print(f"Error reading Word document: {e}")


def extract_text_from_excel(file_path):
    print(f"\n[Extracting text from Spreadsheet]: {file_path}")
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".csv":
            df = pd.read_csv(file_path)
            # Print column headers
            print("| " + " | ".join(df.columns.astype(str)) + " |")
            # Print separator line
            print("|" + "|".join(["---"] * len(df.columns)) + "|")
            # Print each row
            for _, row in df.iterrows():
                print("| " + " | ".join(row.astype(str)) + " |")

        elif ext in [".xlsx", ".xlsm", ".xltx", ".xltm"]:
            wb = load_workbook(file_path, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                print(f"\n--- Sheet: {sheet} ---")
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                # Print header
                print("| " + " | ".join([str(cell) if cell is not None else '' for cell in rows[0]]) + " |")
                # Print separator
                print("|" + "|".join(["---"] * len(rows[0])) + "|")
                # Print rest of the data
                for row in rows[1:]:
                    print("| " + " | ".join([str(cell) if cell is not None else '' for cell in row]) + " |")

        else:
            print("Unsupported spreadsheet format. Only .csv, .xlsx, .xlsm, .xltx, .xltm are supported.")
    except Exception as e:
        print(f"Error reading spreadsheet: {e}")



def extract_text_from_pptx(pptx_path):
    print(f"\n[Extracting text from PowerPoint]: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error reading PowerPoint: {e}")

def extract_text_from_pdf(pdf_path):
    print(f"\n[Extracting text from PDF]: {pdf_path}")
    try:
        doc = fitz.open(pdf_path)
        for page_num, page in enumerate(doc, start=1):
            print(f"\n--- Page {page_num} ---")
            print(page.get_text())
    except Exception as e:
        print(f"Error reading PDF: {e}")


# Example usage
if __name__ == "__main__":
    extract_text_from_word("Apurba Pal resume.docx")
    extract_text_from_excel("HR-Employee-Attrition.csv")
    extract_text_from_pptx("IP-Addressing-and-Masking-A-Foundation-for-Network-Communication.pptx")
    extract_text_from_pdf("Computer Network MCQ.pdf")
